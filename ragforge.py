import os
import io
import json
import threading
import tkinter as tk
from tkinter import filedialog, messagebox, ttk
from datetime import datetime

import cv2
import numpy as np
import xxhash

from PIL import Image
import pytesseract
import whisper

from moviepy.video.io.VideoFileClip import VideoFileClip
from moviepy.audio.io.AudioFileClip import AudioFileClip

import docx
import PyPDF2
import fitz
from pptx import Presentation
from openpyxl import load_workbook
from odf.opendocument import load
from odf.text import P

from transformers import BlipProcessor, BlipForConditionalGeneration

# ===== CONFIG =====

pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
CACHE_FILE = "cache.json"

# ===== GLOBAL MODELS =====

whisper_model = None
blip_processor = None
blip_model = None


def init_models():
    global whisper_model, blip_processor, blip_model

    if whisper_model is None:
        print("[INIT] Loading Whisper...")
        whisper_model = whisper.load_model("base")

    if blip_processor is None:
        print("[INIT] Loading BLIP...")
        blip_processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-base")
        blip_model = BlipForConditionalGeneration.from_pretrained("Salesforce/blip-image-captioning-base")


# ===== CACHE =====

def load_cache():
    if os.path.exists(CACHE_FILE):
        with open(CACHE_FILE, "r") as f:
            return json.load(f)
    return {}


def save_cache(cache):
    with open(CACHE_FILE, "w") as f:
        json.dump(cache, f)

def file_hash(path):
    h = xxhash.xxh64()
    with open(path, "rb") as f:
        while chunk := f.read(8192):
            h.update(chunk)
    return h.hexdigest()


# ===== VIDEO =====

def extract_text_from_video(path):
    try:
        temp_audio = f"temp_{datetime.now().timestamp()}.wav"
        clip = VideoFileClip(path)

        if clip.audio is None:
            return "[NO AUDIO TRACK]"

        try:
            clip.audio.write_audiofile(temp_audio, logger=None)
        except TypeError:
            clip.audio.write_audiofile(temp_audio)

        text = whisper_model.transcribe(temp_audio)["text"]
        os.remove(temp_audio)

        return text

    except Exception as e:
        return f"[VIDEO ERROR] {e}"


def extract_video_frames(path, num_frames=2):
    frames = []
    cap = cv2.VideoCapture(path)

    total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
    if total_frames <= 0:
        return frames

    indices = np.linspace(0, total_frames - 1, num=num_frames, dtype=int)

    for idx in indices:
        cap.set(cv2.CAP_PROP_POS_FRAMES, idx)
        ret, frame = cap.read()
        if ret:
            frames.append(frame)

    cap.release()
    return frames


def describe_frame(frame):
    try:
        image = Image.fromarray(cv2.cvtColor(frame, cv2.COLOR_BGR2RGB))
        inputs = blip_processor(image, return_tensors="pt")
        outputs = blip_model.generate(**inputs)
        return blip_processor.decode(outputs[0], skip_special_tokens=True)
    except Exception as e:
        return f"[FRAME ERROR] {e}"


# ===== IMAGE =====

def generate_image_description(image):
    try:
        inputs = blip_processor(image, return_tensors="pt")
        outputs = blip_model.generate(**inputs)
        return blip_processor.decode(outputs[0], skip_special_tokens=True)
    except Exception as e:
        return f"[BLIP ERROR] {e}"


def analyze_image_bytes(img_bytes):
    try:
        image = Image.open(io.BytesIO(img_bytes)).convert("RGB")

        ocr_text = pytesseract.image_to_string(image)
        caption = generate_image_description(image)

        return f"[Image OCR]\n{ocr_text.strip()}\n\n[Image Caption]\n{caption.strip()}"

    except Exception as e:
        return f"[IMAGE ERROR] {e}"


# ===== DOCUMENT PROCESSING =====

def extract_text_from_pdf(path):
    text = ""

    try:
        with open(path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text += page.extract_text() or ""
    except:
        pass

    try:
        doc = fitz.open(path)
        for page in doc:
            for img in page.get_images(full=True):
                base_image = doc.extract_image(img[0])
                text += "\n\n" + analyze_image_bytes(base_image["image"])
    except:
        pass

    return text


def extract_text_from_docx(path):
    text = ""

    try:
        doc = docx.Document(path)
        for p in doc.paragraphs:
            text += p.text + "\n"
    except:
        pass

    return text


def extract_text_from_pptx(path):
    text = ""

    try:
        pres = Presentation(path)
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except:
        pass

    return text

def extract_text_from_xlsx(path):
    text = ""

    try:
        wb = load_workbook(path, data_only=True)

        for sheet in wb:
            text += f"\n[Sheet: {sheet.title}]\n"
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join([str(cell) for cell in row if cell is not None])
                if row_text:
                    text += row_text + "\n"

    except Exception as e:
        return f"[XLSX ERROR] {e}"

    return text

def extract_text_from_odf(path):
    text = ""

    try:
        doc = load(path)
        paragraphs = doc.getElementsByType(P)

        for p in paragraphs:
            # estrai tutto il testo contenuto nei nodi figli
            for node in p.childNodes:
                if hasattr(node, "data"):
                    text += node.data
            text += "\n"

    except Exception as e:
        return f"[ODF ERROR] {e}"

    return text

# ===== ROUTER =====

def extract_text(path):
    init_models()

    ext = os.path.splitext(path)[1].lower().replace(".", "")

    try:
        if ext in ["jpg", "jpeg", "png"]:
            with open(path, "rb") as f:
                return analyze_image_bytes(f.read())

        elif ext in ["mp3", "wav","m4a"]:
            return whisper_model.transcribe(path)["text"]

        elif ext in ["mp4", "avi", "mov"]:
            text = extract_text_from_video(path)
            for f in extract_video_frames(path):
                text += "\n\n[Frame]\n" + describe_frame(f)
            return text

        elif ext == "pdf":
            return extract_text_from_pdf(path)

        elif ext == "docx":
            return extract_text_from_docx(path)

        elif ext == "pptx":
            return extract_text_from_pptx(path)
        elif ext == "xlsx":
            return extract_text_from_xlsx(path)

        elif ext in ["odt", "ods", "odp"]:
            return extract_text_from_odf(path)

        else:
            # fallback: try raw text read
            try:
                with open(path, "r", encoding="utf-8", errors="ignore") as f:
                    return f.read()
            except:
                return f"[UNSUPPORTED FILE: {ext}]"

    except Exception as e:
        return f"[ERROR PROCESSING FILE: {e}]"


# ===== FILE GROUPS =====

EXT_GROUPS = {
    "Images (jpg, jpeg, png)": ["jpg", "jpeg", "png"],
    "Audio (mp3, wav, m4a)": ["mp3", "wav","m4a"],
    "Video (mp4, avi, mov)": ["mp4", "avi", "mov"],
    "Documents (pdf, docx, pptx, xlsx,odp, odt, ods)": ["pdf", "docx", "pptx","odp","xlsx","ods","odt"]
}


# ===== GUI =====

class App:
    def __init__(self, root):
        self.root = root
        self.root.title("RAGForge by Nanni Bassetti")

        self.input_dir = ""
        self.output_dir = ""
        self.cache = load_cache()

        # INPUT
        self.input_label = tk.Label(root, text="Input: -")
        self.input_label.pack()
        tk.Button(root, text="Select INPUT folder", command=self.select_input).pack()

        # OUTPUT
        self.output_label = tk.Label(root, text="Output: -")
        self.output_label.pack()
        tk.Button(root, text="Select OUTPUT folder", command=self.select_output).pack()

        # ALL FILES
        self.all_files = tk.BooleanVar(value=False)
        tk.Checkbutton(
            root,
            text="ALL (process all files, including unknown formats)",
            variable=self.all_files,
            command=self.toggle_filters
        ).pack()

        # FILTERS
        self.filters = {}
        self.filter_checkbuttons = []

        for key in EXT_GROUPS:
            var = tk.BooleanVar(value=True)
            cb = tk.Checkbutton(root, text=key, variable=var)
            cb.pack(anchor="w")
            self.filters[key] = var
            self.filter_checkbuttons.append(cb)

        # CACHE
        self.use_cache = tk.BooleanVar(value=True)
        tk.Checkbutton(
            root,
            text="Process only new/changed files (use cache)",
            variable=self.use_cache
        ).pack()
        tk.Button(root, text="CLEAR CACHE", command=self.clear_cache).pack(pady=5)
        # COUNT
        self.count_label = tk.Label(root, text="Files to process: 0")
        self.count_label.pack()

        # PROGRESS BAR
        self.progress = ttk.Progressbar(root, length=400, mode="determinate")
        self.progress.pack(pady=10)

        # LOG AREA
        self.log_area = tk.Text(root, height=15, width=80)
        self.log_area.pack()

        # BUTTONS
        tk.Button(root, text="START", command=self.start).pack(pady=5)
        tk.Button(root, text="EXIT", command=root.quit).pack()
    def clear_cache(self):
        confirm = messagebox.askyesno("Confirm", "Delete cache file?")

        if not confirm:
            return

        try:
            if os.path.exists(CACHE_FILE):
                os.remove(CACHE_FILE)

            self.cache = {}
            self.log("Cache cleared")

        except Exception as e:
            self.log(f"Error clearing cache: {e}")
    def log(self, msg):
        self.log_area.insert(tk.END, msg + "\n")
        self.log_area.see(tk.END)
        self.root.update_idletasks()

    def toggle_filters(self):
        state = tk.DISABLED if self.all_files.get() else tk.NORMAL
        for cb in self.filter_checkbuttons:
            cb.config(state=state)
        self.update_file_count()

    def select_input(self):
        self.input_dir = filedialog.askdirectory()
        self.input_label.config(text=f"Input: {self.input_dir}")
        self.update_file_count()

    def select_output(self):
        self.output_dir = filedialog.askdirectory()
        self.output_label.config(text=f"Output: {self.output_dir}")

    def get_active_extensions(self):
        if self.all_files.get():
            return None

        exts = []
        for k, v in self.filters.items():
            if v.get():
                exts.extend(EXT_GROUPS[k])
        return exts

    def update_file_count(self):
        if not self.input_dir:
            return

        exts = self.get_active_extensions()
        count = 0

        for root_dir, _, files in os.walk(self.input_dir):
            for f in files:
                if exts is None:
                    count += 1
                else:
                    ext = os.path.splitext(f)[1].lower().replace(".", "")
                    if ext in exts:
                        count += 1

        self.count_label.config(text=f"Files to process: {count}")

    def start(self):
        if not self.input_dir or not self.output_dir:
            messagebox.showerror("Error", "Select both folders")
            return

        threading.Thread(target=self.run_pipeline).start()

    def run_pipeline(self):
        exts = self.get_active_extensions()

        files = []
        for root_dir, _, filenames in os.walk(self.input_dir):
            for f in filenames:
                full = os.path.join(root_dir, f)

                if exts is None:
                    files.append(full)
                else:
                    ext = os.path.splitext(f)[1].lower().replace(".", "")
                    if ext in exts:
                        files.append(full)

        total = len(files)
        self.progress["maximum"] = total
        self.log(f"Total files: {total}")

        for i, path in enumerate(files, 1):
            filename = os.path.basename(path)
            out_path = os.path.join(self.output_dir, os.path.splitext(filename)[0] + os.path.splitext(filename)[1] + ".txt")

            try:
                h = file_hash(path)

                if self.use_cache.get():
                    if path in self.cache and self.cache[path] == h:
                        self.log(f"SKIP CACHE {filename}")
                        self.progress["value"] = i
                        continue

                text = extract_text(path)

                with open(out_path, "w", encoding="utf-8") as f:
                    f.write(text)

                self.cache[path] = h
                self.log(f"OK {filename}")

            except Exception as e:
                self.log(f"ERROR {filename}: {e}")

            self.progress["value"] = i
            self.root.update_idletasks()

        save_cache(self.cache)
        messagebox.showinfo("Done", "Processing completed!")


# ===== RUN =====

if __name__ == "__main__":
    root = tk.Tk()
    app = App(root)
    root.mainloop()